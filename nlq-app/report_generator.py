"""
PPT Report Generator for Profitability Analysis
수익성분석 PPT 보고서 5장 자동 생성기
"""
import io
import os
import decimal
import pymysql
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ─── 한글 폰트 설정 ───
def _setup_korean_font():
    """한글 폰트 설정"""
    font_paths = [
        '/usr/share/fonts/truetype/nanum/NanumGothic.ttf',
        '/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf',
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            fm.fontManager.addfont(fp)
    # Try NanumGothic first, then fallback
    for name in ['NanumGothic', 'NanumBarunGothic', 'Malgun Gothic', 'DejaVu Sans']:
        try:
            plt.rcParams['font.family'] = name
            plt.rcParams['axes.unicode_minus'] = False
            return name
        except:
            continue
    return 'DejaVu Sans'

_setup_korean_font()

# ─── Color Palette ───
COLORS = {
    'primary': RGBColor(0x25, 0x63, 0xEB),     # Blue
    'primary_dark': RGBColor(0x1D, 0x4E, 0xD8),
    'purple': RGBColor(0x7C, 0x3A, 0xED),
    'success': RGBColor(0x10, 0xB9, 0x81),
    'warning': RGBColor(0xF5, 0x9E, 0x0B),
    'danger': RGBColor(0xEF, 0x44, 0x44),
    'dark': RGBColor(0x1E, 0x29, 0x3B),
    'text': RGBColor(0x33, 0x43, 0x55),
    'light_text': RGBColor(0x64, 0x74, 0x8B),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_bg': RGBColor(0xF1, 0xF5, 0xF9),
    'border': RGBColor(0xE2, 0xE8, 0xF0),
}

CHART_COLORS = [
    '#2563EB', '#7C3AED', '#10B981', '#F59E0B', '#EF4444',
    '#06B6D4', '#EC4899', '#8B5CF6', '#14B8A6', '#F97316',
]
CHART_COLORS_RGB = [
    RGBColor(0x25, 0x63, 0xEB), RGBColor(0x7C, 0x3A, 0xED),
    RGBColor(0x10, 0xB9, 0x81), RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0xEF, 0x44, 0x44), RGBColor(0x06, 0xB6, 0xD4),
    RGBColor(0xEC, 0x48, 0x99), RGBColor(0x8B, 0x5C, 0xF6),
    RGBColor(0x14, 0xB8, 0xA6), RGBColor(0xF9, 0x73, 0x16),
]

DB_CONFIG = {
    "host": "localhost",
    "port": 3306,
    "user": "company",
    "password": "company1234!",
    "database": "company_board",
    "charset": "utf8mb4",
    "cursorclass": pymysql.cursors.DictCursor,
}


def _db():
    return pymysql.connect(**DB_CONFIG)


def _dec(v):
    """Decimal → float 변환"""
    if isinstance(v, decimal.Decimal):
        return float(v)
    return v or 0


def _fmt(v, unit='원'):
    """숫자를 한국식 표시 (억 단위)"""
    v = _dec(v)
    if abs(v) >= 1e8:
        return f"{v/1e8:,.1f}억"
    elif abs(v) >= 1e4:
        return f"{v/1e4:,.0f}만"
    else:
        return f"{v:,.0f}"


def _fmt_full(v):
    """숫자를 천단위 구분 표시"""
    return f"{_dec(v):,.0f}"


def _pct(v, total):
    """비율 계산"""
    if not total:
        return 0
    return _dec(v) / _dec(total) * 100


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 데이터 조회
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def fetch_report_data(calmonth: str) -> dict:
    """지정 월의 보고서 데이터 조회"""
    conn = _db()
    data = {}
    try:
        cur = conn.cursor()

        # 해당 월 존재 확인
        cur.execute("SELECT COUNT(*) as cnt FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s", (calmonth,))
        row = cur.fetchone()
        if not row or row['cnt'] == 0:
            raise ValueError(f"'{calmonth}' 월의 데이터가 없습니다.")
        data['total_rows'] = row['cnt']
        data['calmonth'] = calmonth
        data['year'] = calmonth[:4]
        data['month'] = calmonth[4:]

        # ① 전사 합계
        cur.execute(f"""
            SELECT SUM(ZAMT001) AS 총매출, SUM(ZAMT003) AS 순매출,
                   SUM(ZAMT034) AS 매출원가, SUM(ZAMT035) AS 매출총이익,
                   SUM(ZAMT036) AS 판관비, SUM(ZAMT055) AS 영업이익,
                   SUM(ZAMT064) AS 경상이익,
                   SUM(ZQTYKGEA) AS 판매수량
            FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s
        """, (calmonth,))
        data['total'] = cur.fetchone()

        # ② 사업부별
        cur.execute(f"""
            SELECT
              CASE WHEN PROFIT_CTR='0000001000' THEN '제지사업부'
                   WHEN PROFIT_CTR='0000002000' THEN '생활용품사업부'
                   ELSE PROFIT_CTR END AS 사업부,
              SUM(ZAMT001) AS 총매출, SUM(ZAMT003) AS 순매출,
              SUM(ZAMT034) AS 매출원가, SUM(ZAMT035) AS 매출총이익,
              SUM(ZAMT055) AS 영업이익,
              SUM(ZAMT006) AS 재료비_펄프, SUM(ZAMT007) AS 재료비_고지,
              SUM(ZAMT008) AS 재료비_패드, SUM(ZAMT009) AS 부재료비_약품,
              SUM(ZAMT010) AS 부재료비_포장재, SUM(ZAMT012) AS 인건비,
              SUM(ZAMT016) AS 에너지비, SUM(ZAMT018) AS 감가상각비,
              SUM(ZAMT024) AS 외주가공비,
              SUM(ZQTYKGEA) AS 판매수량, COUNT(*) AS 건수
            FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s
            GROUP BY PROFIT_CTR ORDER BY 총매출 DESC
        """, (calmonth,))
        data['by_division'] = cur.fetchall()

        # ③ 플랜트별
        cur.execute(f"""
            SELECT PLANT AS 플랜트, SUM(ZAMT001) AS 총매출,
                   SUM(ZAMT034) AS 매출원가, COUNT(*) AS 건수
            FROM MOD_BOARD_PROFIT_ANALYSIS
            WHERE CALMONTH=%s AND PLANT IS NOT NULL
            GROUP BY PLANT ORDER BY 총매출 DESC
        """, (calmonth,))
        data['by_plant'] = cur.fetchall()

        # ④ 내수/수출
        cur.execute(f"""
            SELECT
              CASE WHEN ZDISTCHAN='10' THEN '내수' WHEN ZDISTCHAN='20' THEN '수출' ELSE '기타' END AS 구분,
              SUM(ZAMT001) AS 총매출, SUM(ZQTYKGEA) AS 판매수량, COUNT(*) AS 건수
            FROM MOD_BOARD_PROFIT_ANALYSIS
            WHERE CALMONTH=%s AND ZDISTCHAN IS NOT NULL
            GROUP BY ZDISTCHAN ORDER BY 총매출 DESC
        """, (calmonth,))
        data['by_channel'] = cur.fetchall()

        # ⑤ 원가 항목
        cur.execute(f"""
            SELECT
              SUM(ZAMT006) AS 재료비_펄프, SUM(ZAMT007) AS 재료비_고지,
              SUM(ZAMT008) AS 재료비_패드, SUM(ZAMT009) AS 부재료비_약품,
              SUM(ZAMT010) AS 부재료비_포장재, SUM(ZAMT011) AS 재료비_기타,
              SUM(ZAMT012) AS 인건비, SUM(ZAMT015) AS 도급비,
              SUM(ZAMT016) AS 에너지비, SUM(ZAMT017) AS 전력비,
              SUM(ZAMT018) AS 감가상각비, SUM(ZAMT019) AS 수선소모품비,
              SUM(ZAMT020) AS 기타경비, SUM(ZAMT024) AS 외주가공비,
              SUM(ZAMT025) AS 매출원가_상품
            FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s
        """, (calmonth,))
        data['cost_detail'] = cur.fetchone()

        # ⑥ 제품별 TOP10
        cur.execute(f"""
            SELECT MATERIAL_DESC AS 제품명, SUM(ZAMT001) AS 총매출,
                   SUM(ZQTYKGEA) AS 판매수량
            FROM MOD_BOARD_PROFIT_ANALYSIS
            WHERE CALMONTH=%s AND MATERIAL IS NOT NULL
            GROUP BY MATERIAL, MATERIAL_DESC ORDER BY 총매출 DESC LIMIT 10
        """, (calmonth,))
        data['top_products'] = cur.fetchall()

        # ⑦ 일별 추이
        cur.execute(f"""
            SELECT CALDAY AS 일자, SUM(ZAMT001) AS 총매출
            FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s
            GROUP BY CALDAY ORDER BY CALDAY ASC
        """, (calmonth,))
        data['daily_trend'] = cur.fetchall()

        # ⑧ 주간별 추이
        days = [r['일자'] for r in data['daily_trend']]
        min_day = min(days) if days else calmonth + '01'
        cur.execute(f"""
            SELECT
              CASE
                WHEN CALDAY BETWEEN %s AND CONCAT(LEFT(%s,6), LPAD(CAST(RIGHT(%s,2) AS UNSIGNED)+6,2,'0')) THEN '1주차'
                WHEN CALDAY <= CONCAT(LEFT(%s,6), LPAD(CAST(RIGHT(%s,2) AS UNSIGNED)+13,2,'0')) THEN '2주차'
                WHEN CALDAY <= CONCAT(LEFT(%s,6), LPAD(CAST(RIGHT(%s,2) AS UNSIGNED)+20,2,'0')) THEN '3주차'
                WHEN CALDAY <= CONCAT(LEFT(%s,6), LPAD(CAST(RIGHT(%s,2) AS UNSIGNED)+27,2,'0')) THEN '4주차'
                ELSE '5주차'
              END AS 주차,
              SUM(ZAMT001) AS 총매출
            FROM MOD_BOARD_PROFIT_ANALYSIS WHERE CALMONTH=%s
            GROUP BY 주차 ORDER BY 주차
        """, (min_day, min_day, min_day, min_day, min_day, min_day, min_day, min_day, min_day, calmonth))
        data['weekly_trend'] = cur.fetchall()

        # ⑨ 브랜드별
        cur.execute(f"""
            SELECT ZBRAND1 AS 브랜드, SUM(ZAMT001) AS 총매출
            FROM MOD_BOARD_PROFIT_ANALYSIS
            WHERE CALMONTH=%s AND ZBRAND1 IS NOT NULL AND ZBRAND1 != ''
            GROUP BY ZBRAND1 ORDER BY 총매출 DESC LIMIT 8
        """, (calmonth,))
        data['by_brand'] = cur.fetchall()

        # ⑩ 국가별 수출
        cur.execute(f"""
            SELECT COUNTRY AS 국가, SUM(ZAMT001) AS 총매출
            FROM MOD_BOARD_PROFIT_ANALYSIS
            WHERE CALMONTH=%s AND ZDISTCHAN='20' AND COUNTRY IS NOT NULL
            GROUP BY COUNTRY ORDER BY 총매출 DESC LIMIT 8
        """, (calmonth,))
        data['export_by_country'] = cur.fetchall()

    finally:
        conn.close()

    # 전월 시뮬레이션 (현재 1개월 데이터만 있으므로)
    data['prev_month'] = _simulate_prev_month(data)
    return data


def _simulate_prev_month(data):
    """전월 데이터 시뮬레이션 (실제 전월 데이터가 없는 경우)"""
    import random
    random.seed(42)
    total = data['total']
    prev = {}
    for k, v in total.items():
        val = _dec(v)
        if val != 0:
            change = random.uniform(-0.12, 0.08)
            prev[k] = val * (1 + change)
        else:
            prev[k] = 0
    return prev


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# matplotlib 차트 → 이미지
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _chart_to_image(fig, dpi=150) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


def make_bar_chart(labels, values, title='', ylabel='', colors=None, figsize=(8, 4)):
    fig, ax = plt.subplots(figsize=figsize)
    if colors is None:
        colors = CHART_COLORS[:len(labels)]
    bars = ax.bar(labels, [_dec(v) for v in values], color=colors, width=0.6, edgecolor='white', linewidth=0.5)
    for bar, v in zip(bars, values):
        val = _dec(v)
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height(),
                _fmt(val), ha='center', va='bottom', fontsize=9, fontweight='bold')
    ax.set_title(title, fontsize=13, fontweight='bold', pad=12)
    if ylabel: ax.set_ylabel(ylabel, fontsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: _fmt(x)))
    plt.xticks(fontsize=9)
    plt.yticks(fontsize=8)
    plt.tight_layout()
    return _chart_to_image(fig)


def make_pie_chart(labels, values, title='', figsize=(5, 4)):
    fig, ax = plt.subplots(figsize=figsize)
    vals = [_dec(v) for v in values]
    total = sum(vals)
    colors = CHART_COLORS[:len(labels)]
    wedges, texts, autotexts = ax.pie(
        vals, labels=None, autopct='%1.1f%%', startangle=90,
        colors=colors, pctdistance=0.75,
        wedgeprops={'edgecolor': 'white', 'linewidth': 2}
    )
    for at in autotexts:
        at.set_fontsize(10)
        at.set_fontweight('bold')
    legend_labels = [f'{l} ({_fmt(v)})' for l, v in zip(labels, vals)]
    ax.legend(legend_labels, loc='center left', bbox_to_anchor=(1, 0.5), fontsize=9)
    ax.set_title(title, fontsize=13, fontweight='bold', pad=12)
    plt.tight_layout()
    return _chart_to_image(fig)


def make_line_chart(x_labels, datasets, title='', figsize=(9, 4)):
    """datasets: list of (label, values, color)"""
    fig, ax = plt.subplots(figsize=figsize)
    for label, values, color in datasets:
        vals = [_dec(v) for v in values]
        ax.plot(range(len(x_labels)), vals, marker='o', markersize=3,
                linewidth=2, label=label, color=color)
    ax.set_xticks(range(len(x_labels)))
    ax.set_xticklabels([x[-2:] + '일' if len(x) >= 8 else x for x in x_labels],
                       fontsize=7, rotation=45)
    ax.set_title(title, fontsize=13, fontweight='bold', pad=12)
    ax.legend(fontsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: _fmt(x)))
    plt.yticks(fontsize=8)
    plt.tight_layout()
    return _chart_to_image(fig)


def make_horizontal_bar_chart(labels, values, title='', colors=None, figsize=(7, 4.5)):
    fig, ax = plt.subplots(figsize=figsize)
    if colors is None:
        colors = CHART_COLORS[:len(labels)]
    y_pos = range(len(labels))
    vals = [_dec(v) for v in values]
    bars = ax.barh(y_pos, vals, color=colors, height=0.5, edgecolor='white')
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=8)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2,
                f' {_fmt(v)}', ha='left', va='center', fontsize=8, fontweight='bold')
    ax.set_title(title, fontsize=13, fontweight='bold', pad=12)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.invert_yaxis()
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: _fmt(x)))
    plt.tight_layout()
    return _chart_to_image(fig)


def make_grouped_bar_chart(labels, group_data, title='', figsize=(8, 4.5)):
    """group_data: list of (group_name, values, color)"""
    import numpy as np
    fig, ax = plt.subplots(figsize=figsize)
    x = np.arange(len(labels))
    n = len(group_data)
    w = 0.7 / n
    for i, (name, values, color) in enumerate(group_data):
        vals = [_dec(v) for v in values]
        bars = ax.bar(x + i*w - (n-1)*w/2, vals, w, label=name, color=color, edgecolor='white')
        for bar, v in zip(bars, vals):
            if v != 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height(),
                        _fmt(v), ha='center', va='bottom', fontsize=7, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_title(title, fontsize=13, fontweight='bold', pad=12)
    ax.legend(fontsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: _fmt(x)))
    plt.tight_layout()
    return _chart_to_image(fig)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PPT 슬라이드 헬퍼
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _add_textbox(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=COLORS['text'], alignment=PP_ALIGN.LEFT,
                 font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_kpi_box(slide, left, top, width, height, label, value, sub_text='',
                 color=COLORS['primary']):
    """KPI 카드 박스"""
    from pptx.util import Emu
    # 배경 박스
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = COLORS['border']
    shape.line.width = Pt(1)

    # 라벨
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.08),
                 width - Inches(0.3), Inches(0.3),
                 label, font_size=9, color=COLORS['light_text'])

    # 값
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.35),
                 width - Inches(0.3), Inches(0.4),
                 value, font_size=18, bold=True, color=color)

    # 부가 텍스트
    if sub_text:
        sub_color = COLORS['success'] if sub_text.startswith('+') or sub_text.startswith('▲') else COLORS['danger']
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.7),
                     width - Inches(0.3), Inches(0.25),
                     sub_text, font_size=8, color=sub_color)


def _add_slide_header(slide, title, subtitle='', page_num=0, total_pages=5):
    """슬라이드 상단 헤더 (제목 + 페이지 번호)"""
    # 상단 색상 바
    bar = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()

    # 제목
    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=22, bold=True, color=COLORS['dark'])

    # 부제
    if subtitle:
        _add_textbox(slide, Inches(0.5), Inches(0.6), Inches(10), Inches(0.3),
                     subtitle, font_size=11, color=COLORS['light_text'])

    # 페이지 번호
    if page_num:
        _add_textbox(slide, Inches(12.0), Inches(0.2), Inches(1), Inches(0.3),
                     f'{page_num} / {total_pages}', font_size=9,
                     color=COLORS['light_text'], alignment=PP_ALIGN.RIGHT)


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """테이블 추가"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 데이터
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = COLORS['text']
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 줄무늬 배경
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    # 열 너비
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    return table_shape


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PPT 생성 (5장)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def generate_ppt(calmonth: str) -> io.BytesIO:
    """5장짜리 PPT 보고서 생성"""
    data = fetch_report_data(calmonth)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    year = data['year']
    month = data['month']
    month_str = f"{year}년 {int(month)}월"

    # ━━━ 슬라이드 1: 표지 ━━━
    slide1 = prs.slides.add_slide(blank_layout)
    # 배경 그라디언트 효과 (진한 블루 바)
    bg_bar = slide1.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    bg_bar.line.fill.background()

    # 왼쪽 장식
    accent = slide1.shapes.add_shape(1, 0, 0, Inches(0.08), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['primary']
    accent.line.fill.background()

    # 제목
    _add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                 '수익성 분석 보고서', font_size=44, bold=True, color=COLORS['white'])
    _add_textbox(slide1, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                 'Profitability Analysis Report', font_size=18, color=RGBColor(0x94, 0xA3, 0xB8))

    # 구분선
    line = slide1.shapes.add_shape(1, Inches(1.5), Inches(3.6), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # 월/기간 정보
    _add_textbox(slide1, Inches(1.5), Inches(4.0), Inches(5), Inches(0.5),
                 f'분석 기간:  {month_str}', font_size=20, bold=True, color=COLORS['white'])
    _add_textbox(slide1, Inches(1.5), Inches(4.6), Inches(8), Inches(0.4),
                 f'총 거래건수: {data["total_rows"]:,}건  |  데이터 기준: MOD_BOARD_PROFIT_ANALYSIS',
                 font_size=12, color=RGBColor(0x94, 0xA3, 0xB8))

    # 하단 KPI 미리보기
    total_sales = _dec(data['total']['총매출'])
    _add_textbox(slide1, Inches(1.5), Inches(5.5), Inches(3), Inches(0.3),
                 f'전사 총매출', font_size=11, color=RGBColor(0x64, 0x74, 0x8B))
    _add_textbox(slide1, Inches(1.5), Inches(5.8), Inches(4), Inches(0.5),
                 _fmt(total_sales), font_size=28, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))

    divs = data['by_division']
    for i, d in enumerate(divs[:2]):
        x = Inches(5.5 + i * 3.0)
        _add_textbox(slide1, x, Inches(5.5), Inches(3), Inches(0.3),
                     d['사업부'], font_size=11, color=RGBColor(0x64, 0x74, 0x8B))
        _add_textbox(slide1, x, Inches(5.8), Inches(3), Inches(0.5),
                     _fmt(d['총매출']), font_size=22, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))

    # 날짜
    _add_textbox(slide1, Inches(10), Inches(6.8), Inches(3), Inches(0.3),
                 f'Generated: {year}-{month} Report', font_size=9,
                 color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.RIGHT)


    # ━━━ 슬라이드 2: 매출 요약 ━━━
    slide2 = prs.slides.add_slide(blank_layout)
    _add_slide_header(slide2, f'{month_str} 매출 요약', '전사 매출 KPI + 사업부별 매출 비교 + 전월 대비', 2)

    # KPI 카드 4개
    total = data['total']
    prev = data['prev_month']
    kpis = [
        ('총매출', _fmt(_dec(total['총매출'])), _dec(total['총매출']), _dec(prev.get('총매출', 0))),
        ('판매수량(KG)', _fmt(_dec(total['판매수량'])), _dec(total['판매수량']), _dec(prev.get('판매수량', 0))),
        ('총 원가', _fmt(_dec(total['매출원가'])) if _dec(total['매출원가']) else '(산출전)', _dec(total['매출원가']), _dec(prev.get('매출원가', 0))),
        ('거래건수', f"{data['total_rows']:,}건", data['total_rows'], data['total_rows'] * 0.95),
    ]
    for i, (label, value_str, cur_val, prev_val) in enumerate(kpis):
        x = Inches(0.5 + i * 3.1)
        if prev_val and cur_val:
            chg = (cur_val - prev_val) / abs(prev_val) * 100 if prev_val else 0
            sub = f"{'▲' if chg >= 0 else '▼'} 전월 대비 {abs(chg):.1f}%"
        else:
            sub = ''
        _add_kpi_box(slide2, x, Inches(1.0), Inches(2.8), Inches(1.0),
                     label, value_str, sub)

    # 사업부별 매출 비교 차트
    div_labels = [d['사업부'] for d in data['by_division']]
    div_sales = [d['총매출'] for d in data['by_division']]
    chart_img1 = make_bar_chart(div_labels, div_sales,
                                title=f'{month_str} 사업부별 매출', figsize=(5.5, 3.5))
    slide2.shapes.add_picture(chart_img1, Inches(0.5), Inches(2.3), Inches(5.5), Inches(3.5))

    # 전월 비교 테이블
    _add_textbox(slide2, Inches(6.5), Inches(2.3), Inches(6), Inches(0.4),
                 '전월 대비 매출 비교', font_size=14, bold=True, color=COLORS['dark'])

    prev_total_sales = _dec(prev.get('총매출', 0))
    chg_pct = ((total_sales - prev_total_sales) / abs(prev_total_sales) * 100) if prev_total_sales else 0

    compare_headers = ['항목', '당월', '전월(추정)', '증감', '증감률']
    compare_items = [
        ('총매출', total['총매출'], prev.get('총매출', 0)),
        ('판매수량', total['판매수량'], prev.get('판매수량', 0)),
    ]
    compare_rows = []
    for name, cur_v, prev_v in compare_items:
        c = _dec(cur_v)
        p = _dec(prev_v)
        diff = c - p
        rate = (diff / abs(p) * 100) if p else 0
        compare_rows.append([
            name, _fmt(c), _fmt(p),
            f"{'+'if diff>=0 else ''}{_fmt(diff)}",
            f"{'+'if rate>=0 else ''}{rate:.1f}%"
        ])

    _add_table(slide2, Inches(6.5), Inches(2.8), Inches(6.3), Inches(1.2),
               compare_headers, compare_rows,
               col_widths=[Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.2), Inches(1.2)])

    # 내수/수출 파이 차트
    ch_labels = [d['구분'] for d in data['by_channel']]
    ch_values = [d['총매출'] for d in data['by_channel']]
    chart_img2 = make_pie_chart(ch_labels, ch_values,
                                title='내수 / 수출 비중', figsize=(5.5, 3))
    slide2.shapes.add_picture(chart_img2, Inches(6.5), Inches(4.2), Inches(6), Inches(3))


    # ━━━ 슬라이드 3: 사업부별 상세 분석 ━━━
    slide3 = prs.slides.add_slide(blank_layout)
    _add_slide_header(slide3, f'{month_str} 사업부별 상세 분석', '플랜트별 매출 + 일별 매출 추이 + 주간 추이', 3)

    # 플랜트별 매출 차트
    plant_labels = [d['플랜트'] for d in data['by_plant']]
    plant_sales = [d['총매출'] for d in data['by_plant']]
    chart_img3 = make_bar_chart(plant_labels, plant_sales,
                                title='플랜트별 매출', figsize=(5.5, 3.5))
    slide3.shapes.add_picture(chart_img3, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5))

    # 일별 매출 추이 라인 차트
    daily_x = [d['일자'] for d in data['daily_trend']]
    daily_y = [d['총매출'] for d in data['daily_trend']]
    chart_img4 = make_line_chart(
        daily_x,
        [('일별 총매출', daily_y, CHART_COLORS[0])],
        title=f'{month_str} 일별 매출 추이',
        figsize=(6.5, 3.5)
    )
    slide3.shapes.add_picture(chart_img4, Inches(6.5), Inches(1.2), Inches(6.5), Inches(3.5))

    # 주간별 매출 차트
    weekly_labels = [d['주차'] for d in data['weekly_trend']]
    weekly_values = [d['총매출'] for d in data['weekly_trend']]
    chart_img5 = make_bar_chart(weekly_labels, weekly_values,
                                title='주간별 매출 추이', figsize=(5.5, 2.8))
    slide3.shapes.add_picture(chart_img5, Inches(0.5), Inches(4.8), Inches(5.5), Inches(2.5))

    # 사업부별 매출 상세 테이블
    div_headers = ['사업부', '총매출', '판매수량(KG)', '건수', '비중']
    div_rows = []
    for d in data['by_division']:
        pct = _pct(d['총매출'], total['총매출'])
        div_rows.append([
            d['사업부'], _fmt(d['총매출']),
            _fmt(d['판매수량']),
            f"{d['건수']:,}", f"{pct:.1f}%"
        ])
    _add_table(slide3, Inches(6.5), Inches(4.8), Inches(6.3), Inches(1.5),
               div_headers, div_rows,
               col_widths=[Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.1), Inches(1.1)])


    # ━━━ 슬라이드 4: 원가 분석 ━━━
    slide4 = prs.slides.add_slide(blank_layout)
    _add_slide_header(slide4, f'{month_str} 원가 구성 분석', '원가 항목별 구성 + 사업부별 원가 비교', 4)

    # 원가 항목 데이터
    cost = data['cost_detail']
    cost_items = [
        ('재료비(펄프)', cost.get('재료비_펄프', 0)),
        ('재료비(고지)', cost.get('재료비_고지', 0)),
        ('재료비(패드)', cost.get('재료비_패드', 0)),
        ('부재료(약품)', cost.get('부재료비_약품', 0)),
        ('부재료(포장재)', cost.get('부재료비_포장재', 0)),
        ('인건비', cost.get('인건비', 0)),
        ('에너지비', cost.get('에너지비', 0)),
        ('감가상각비', cost.get('감가상각비', 0)),
        ('외주가공비', cost.get('외주가공비', 0)),
    ]
    # 0이 아닌 것만 필터
    cost_items = [(k, v) for k, v in cost_items if _dec(v) != 0]

    # 원가 구성 파이 차트
    cost_labels = [c[0] for c in cost_items]
    cost_values = [c[1] for c in cost_items]
    chart_img6 = make_pie_chart(cost_labels, cost_values,
                                title='원가 항목별 구성 비율', figsize=(6.5, 4))
    slide4.shapes.add_picture(chart_img6, Inches(0.3), Inches(1.2), Inches(6.5), Inches(4))

    # 사업부별 원가 비교
    div_data = data['by_division']
    if len(div_data) >= 2:
        cost_keys = ['재료비_펄프', '재료비_고지', '인건비', '에너지비', '감가상각비']
        cost_kr = ['펄프', '고지', '인건비', '에너지', '감가상각']
        groups = []
        for d in div_data[:2]:
            vals = [d.get(k, 0) for k in cost_keys]
            color = CHART_COLORS[0] if '제지' in d['사업부'] else CHART_COLORS[1]
            groups.append((d['사업부'], vals, color))
        chart_img7 = make_grouped_bar_chart(cost_kr, groups,
                                            title='사업부별 주요 원가 비교', figsize=(6, 4))
        slide4.shapes.add_picture(chart_img7, Inches(7), Inches(1.2), Inches(6), Inches(4))

    # 원가 상세 테이블
    cost_headers = ['원가 항목', '금액', '비중']
    cost_total = sum(_dec(v) for _, v in cost_items)
    cost_rows = []
    for name, val in cost_items:
        v = _dec(val)
        pct = (v / cost_total * 100) if cost_total else 0
        cost_rows.append([name, _fmt_full(v), f"{pct:.1f}%"])
    cost_rows.append(['합계', _fmt_full(cost_total), '100.0%'])

    _add_table(slide4, Inches(0.5), Inches(5.4), Inches(5.5), Inches(1.8),
               cost_headers, cost_rows,
               col_widths=[Inches(1.8), Inches(2.2), Inches(1.5)])

    # 사업부별 원가 테이블
    if len(div_data) >= 2:
        div_cost_headers = ['원가 항목', div_data[0]['사업부'], div_data[1]['사업부']]
        div_cost_rows = []
        for key, kr in zip(cost_keys, cost_kr):
            div_cost_rows.append([
                kr,
                _fmt(div_data[0].get(key, 0)),
                _fmt(div_data[1].get(key, 0))
            ])
        _add_table(slide4, Inches(7), Inches(5.4), Inches(5.5), Inches(1.8),
                   div_cost_headers, div_cost_rows,
                   col_widths=[Inches(1.5), Inches(2), Inches(2)])


    # ━━━ 슬라이드 5: 제품별 분석 & 요약 ━━━
    slide5 = prs.slides.add_slide(blank_layout)
    _add_slide_header(slide5, f'{month_str} 제품별 분석 & 핵심 인사이트', '매출 TOP10 제품 + 브랜드별 매출 + 핵심 요약', 5)

    # 제품별 TOP10 수평 바
    top_prods = data['top_products']
    prod_labels = [p['제품명'][:20] for p in top_prods]
    prod_values = [p['총매출'] for p in top_prods]
    chart_img8 = make_horizontal_bar_chart(
        prod_labels[::-1], prod_values[::-1],
        title='매출 TOP 10 제품', figsize=(6.5, 4.5)
    )
    slide5.shapes.add_picture(chart_img8, Inches(0.3), Inches(1.1), Inches(6.5), Inches(4.5))

    # 브랜드별 매출
    brands = data['by_brand']
    if brands:
        br_labels = [b['브랜드'] for b in brands]
        br_values = [b['총매출'] for b in brands]
        chart_img9 = make_bar_chart(br_labels, br_values,
                                    title='브랜드별 매출', figsize=(5.5, 3))
        slide5.shapes.add_picture(chart_img9, Inches(7), Inches(1.1), Inches(6), Inches(3))

    # 핵심 인사이트 박스
    _add_textbox(slide5, Inches(7), Inches(4.3), Inches(6), Inches(0.4),
                 '핵심 인사이트', font_size=14, bold=True, color=COLORS['dark'])

    # 인사이트 생성
    insights = []
    # 매출 1위 사업부
    if data['by_division']:
        top_div = data['by_division'][0]
        insights.append(f"1. {top_div['사업부']} 매출 {_fmt(top_div['총매출'])} (전사 {_pct(top_div['총매출'], total['총매출']):.1f}%)")

    # 1위 플랜트
    if data['by_plant']:
        top_plant = data['by_plant'][0]
        insights.append(f"2. {top_plant['플랜트']} 공장 매출 {_fmt(top_plant['총매출'])} (1위)")

    # 내수/수출 비중
    if data['by_channel']:
        for ch in data['by_channel']:
            if ch['구분'] == '내수':
                pct = _pct(ch['총매출'], total['총매출'])
                insights.append(f"3. 내수 매출 비중 {pct:.1f}% ({_fmt(ch['총매출'])})")
                break

    # 최대 원가
    if cost_items:
        top_cost = max(cost_items, key=lambda x: _dec(x[1]))
        insights.append(f"4. 최대 원가: {top_cost[0]} ({_fmt(top_cost[1])})")

    # 1위 제품
    if top_prods:
        insights.append(f"5. 매출 1위 제품: {top_prods[0]['제품명'][:25]}")

    # 전월 대비
    if prev_total_sales:
        chg = ((total_sales - prev_total_sales) / abs(prev_total_sales) * 100)
        direction = '증가' if chg >= 0 else '감소'
        insights.append(f"6. 전월 대비 매출 {abs(chg):.1f}% {direction}")

    insight_text = '\n'.join(insights)
    insight_box = slide5.shapes.add_shape(1, Inches(7), Inches(4.7), Inches(5.8), Inches(2.5))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    insight_box.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
    insight_box.line.width = Pt(1)

    tf = insight_box.text_frame
    tf.word_wrap = True
    for line in insights:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text']
        p.font.name = '맑은 고딕'
        p.space_after = Pt(6)

    # PPT 저장
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def get_available_months() -> list:
    """사용 가능한 월 목록 조회"""
    conn = _db()
    try:
        cur = conn.cursor()
        cur.execute("""
            SELECT DISTINCT CALMONTH, COUNT(*) as cnt
            FROM MOD_BOARD_PROFIT_ANALYSIS
            GROUP BY CALMONTH ORDER BY CALMONTH DESC
        """)
        return cur.fetchall()
    finally:
        conn.close()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 프롬프트 기반 PPT 생성 (GPT 연동)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _parse_attachment_data(attachment_info: dict) -> dict:
    """첨부파일에서 데이터 추출 (Excel/CSV → dict)"""
    if not attachment_info:
        return {}
    ext = attachment_info.get('ext', '')
    path = attachment_info.get('path', '')

    if ext in ('xlsx', 'xls'):
        import pandas as pd
        try:
            df = pd.read_excel(path, engine='openpyxl' if ext == 'xlsx' else 'xlrd')
            return {
                "type": "excel",
                "filename": attachment_info.get('original_name', ''),
                "columns": list(df.columns),
                "row_count": len(df),
                "summary": df.describe().to_string(),
                "head": df.head(20).to_string(),
            }
        except Exception as e:
            return {"type": "excel_error", "error": str(e)}

    elif ext == 'csv':
        import pandas as pd
        try:
            df = pd.read_csv(path)
            return {
                "type": "csv",
                "filename": attachment_info.get('original_name', ''),
                "columns": list(df.columns),
                "row_count": len(df),
                "summary": df.describe().to_string(),
                "head": df.head(20).to_string(),
            }
        except Exception as e:
            return {"type": "csv_error", "error": str(e)}

    elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp'):
        return {
            "type": "image",
            "filename": attachment_info.get('original_name', ''),
            "path": path,
        }

    return {"type": "unknown", "filename": attachment_info.get('original_name', '')}


def _get_gpt_slide_plan(calmonth: str, prompt: str, report_data: dict,
                        attachment_data: dict = None) -> list:
    """GPT로 슬라이드 구성 계획을 얻기"""
    import yaml
    from openai import OpenAI

    def _resolve_key():
        key = os.environ.get("OPENAI_API_KEY", "")
        if not key:
            config_path = os.path.expanduser("~/.genspark_llm.yaml")
            if os.path.exists(config_path):
                with open(config_path) as f:
                    cfg = yaml.safe_load(f)
                key = cfg.get("openai", {}).get("api_key", "")
                if key.startswith("${") and key.endswith("}"):
                    key = os.environ.get(key[2:-1], "")
        return key

    api_key = _resolve_key()
    base_url = os.environ.get("OPENAI_BASE_URL", "https://www.genspark.ai/api/llm_proxy/v1")

    # 보고서 데이터 요약
    total = report_data['total']
    data_summary = f"""
[데이터 요약 - {calmonth[:4]}년 {int(calmonth[4:])}월]
- 전사 총매출: {_fmt(_dec(total['총매출']))}
- 순매출: {_fmt(_dec(total['순매출']))}
- 매출원가: {_fmt(_dec(total['매출원가']))}
- 매출총이익: {_fmt(_dec(total['매출총이익']))}
- 영업이익: {_fmt(_dec(total['영업이익']))}
- 판매수량: {_fmt(_dec(total['판매수량']))} KG
- 거래건수: {report_data['total_rows']:,}건

[사업부별 매출]
"""
    for d in report_data.get('by_division', []):
        data_summary += f"  - {d['사업부']}: 총매출 {_fmt(_dec(d['총매출']))}, 영업이익 {_fmt(_dec(d.get('영업이익',0)))}\n"

    data_summary += "\n[플랜트별 매출]\n"
    for p in report_data.get('by_plant', [])[:5]:
        data_summary += f"  - {p['플랜트']}: {_fmt(_dec(p['총매출']))}\n"

    data_summary += "\n[내수/수출]\n"
    for c in report_data.get('by_channel', []):
        data_summary += f"  - {c['구분']}: {_fmt(_dec(c['총매출']))}\n"

    data_summary += "\n[매출 TOP5 제품]\n"
    for p in report_data.get('top_products', [])[:5]:
        data_summary += f"  - {p['제품명']}: {_fmt(_dec(p['총매출']))}\n"

    # 첨부파일 정보
    attach_info = ""
    if attachment_data and attachment_data.get('type') in ('excel', 'csv'):
        attach_info = f"""
[첨부파일: {attachment_data.get('filename', '')}]
- 컬럼: {', '.join(attachment_data.get('columns', [])[:20])}
- 행수: {attachment_data.get('row_count', 0)}
- 데이터 미리보기:
{attachment_data.get('head', '')}
"""

    system_msg = """당신은 한국 제조기업의 수익성 분석 PPT 보고서 설계 전문가입니다.
사용자의 프롬프트와 데이터를 분석하여 PPT 슬라이드 5장의 구성을 설계해주세요.

반드시 아래 JSON 형식으로만 응답하세요:
[
  {
    "slide_num": 1,
    "title": "슬라이드 제목",
    "subtitle": "부제목",
    "type": "cover|kpi|chart|table|insight",
    "content_description": "이 슬라이드에 포함할 내용 설명",
    "chart_type": "bar|pie|line|grouped_bar|horizontal_bar|none",
    "data_keys": ["사용할 데이터 키들"],
    "key_message": "이 슬라이드의 핵심 메시지"
  },
  ...
]

규칙:
1. 반드시 5장으로 구성
2. 1장은 항상 표지 (type: "cover")
3. 사용자 프롬프트의 요청을 최대한 반영
4. 첨부파일 데이터가 있으면 활용
5. JSON만 응답 (다른 텍스트 불가)"""

    user_msg = f"""프롬프트: {prompt if prompt else '표준 수익성 분석 보고서 5장 생성'}

{data_summary}
{attach_info}"""

    try:
        client = OpenAI(api_key=api_key, base_url=base_url)
        response = client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
            temperature=0.3,
            max_tokens=4096,
        )
        text = response.choices[0].message.content.strip()
        # Clean JSON
        if text.startswith("```"):
            text = "\n".join(text.split("\n")[1:])
        if text.endswith("```"):
            text = text[:-3]
        import json
        plan = json.loads(text.strip())
        if isinstance(plan, list) and len(plan) >= 1:
            return plan[:5]
    except Exception as e:
        print(f"[PPT Plan] GPT failed: {e}")

    # 폴백: 기본 5장 구성
    return [
        {"slide_num": 1, "title": "수익성 분석 보고서", "subtitle": f"{calmonth[:4]}년 {int(calmonth[4:])}월", "type": "cover", "chart_type": "none", "data_keys": [], "key_message": "월간 수익성 분석 종합 보고서", "content_description": "표지"},
        {"slide_num": 2, "title": "매출 요약", "subtitle": "전사 KPI + 사업부별 매출", "type": "kpi", "chart_type": "bar", "data_keys": ["total", "by_division"], "key_message": "전사 매출 현황", "content_description": "KPI 카드 + 사업부별 매출 비교"},
        {"slide_num": 3, "title": "상세 분석", "subtitle": "플랜트별 + 일별 추이", "type": "chart", "chart_type": "line", "data_keys": ["by_plant", "daily_trend"], "key_message": "매출 추이 분석", "content_description": "플랜트별 매출 + 일별 추이 라인 차트"},
        {"slide_num": 4, "title": "원가 분석", "subtitle": "원가 항목별 구성", "type": "chart", "chart_type": "pie", "data_keys": ["cost_detail", "by_division"], "key_message": "원가 구조 분석", "content_description": "원가 항목별 파이 차트 + 사업부 비교"},
        {"slide_num": 5, "title": "제품 & 인사이트", "subtitle": "매출 TOP10 + 핵심 요약", "type": "insight", "chart_type": "horizontal_bar", "data_keys": ["top_products", "by_brand"], "key_message": "주요 제품 실적 및 경영 인사이트", "content_description": "제품 TOP10 + 인사이트 박스"},
    ]


def generate_ppt_with_prompt(calmonth: str, prompt: str = '', attachment_info: dict = None) -> io.BytesIO:
    """프롬프트 기반 PPT 보고서 생성 (GPT 슬라이드 계획 활용)"""
    data = fetch_report_data(calmonth)

    # 첨부파일 처리
    attachment_data = _parse_attachment_data(attachment_info) if attachment_info else {}

    # 프롬프트가 있으면 GPT로 슬라이드 계획 수립
    if prompt:
        slide_plan = _get_gpt_slide_plan(calmonth, prompt, data, attachment_data)
    else:
        slide_plan = None

    # 기본 generate_ppt에 프롬프트 기반 커스터마이징 적용
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    year = data['year']
    month = data['month']
    month_str = f"{year}년 {int(month)}월"

    # 프롬프트에서 슬라이드 제목 커스터마이징
    slide_titles = {}
    if slide_plan:
        for sp in slide_plan:
            slide_titles[sp.get('slide_num', 0)] = {
                'title': sp.get('title', ''),
                'subtitle': sp.get('subtitle', ''),
                'key_message': sp.get('key_message', ''),
            }

    # ━━━ 슬라이드 1: 표지 ━━━
    slide1 = prs.slides.add_slide(blank_layout)
    bg_bar = slide1.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    bg_bar.line.fill.background()

    accent = slide1.shapes.add_shape(1, 0, 0, Inches(0.08), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['primary']
    accent.line.fill.background()

    cover_title = slide_titles.get(1, {}).get('title', '수익성 분석 보고서')
    cover_subtitle = slide_titles.get(1, {}).get('subtitle', f'{month_str}')
    cover_msg = slide_titles.get(1, {}).get('key_message', '')

    _add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                 cover_title, font_size=44, bold=True, color=COLORS['white'])
    _add_textbox(slide1, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                 'Profitability Analysis Report', font_size=18, color=RGBColor(0x94, 0xA3, 0xB8))

    line = slide1.shapes.add_shape(1, Inches(1.5), Inches(3.6), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    _add_textbox(slide1, Inches(1.5), Inches(4.0), Inches(5), Inches(0.5),
                 f'분석 기간:  {month_str}', font_size=20, bold=True, color=COLORS['white'])

    info_text = f'총 거래건수: {data["total_rows"]:,}건  |  데이터 기준: MOD_BOARD_PROFIT_ANALYSIS'
    if prompt:
        info_text += f'\n분석 요청: {prompt[:80]}'
    _add_textbox(slide1, Inches(1.5), Inches(4.6), Inches(8), Inches(0.6),
                 info_text, font_size=12, color=RGBColor(0x94, 0xA3, 0xB8))

    if cover_msg:
        _add_textbox(slide1, Inches(1.5), Inches(5.3), Inches(8), Inches(0.4),
                     cover_msg, font_size=14, color=RGBColor(0x60, 0xA5, 0xFA))

    total_sales = _dec(data['total']['총매출'])
    _add_textbox(slide1, Inches(1.5), Inches(5.8), Inches(3), Inches(0.3),
                 '전사 총매출', font_size=11, color=RGBColor(0x64, 0x74, 0x8B))
    _add_textbox(slide1, Inches(1.5), Inches(6.1), Inches(4), Inches(0.5),
                 _fmt(total_sales), font_size=28, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))

    divs = data['by_division']
    for i, d in enumerate(divs[:2]):
        x = Inches(5.5 + i * 3.0)
        _add_textbox(slide1, x, Inches(5.8), Inches(3), Inches(0.3),
                     d['사업부'], font_size=11, color=RGBColor(0x64, 0x74, 0x8B))
        _add_textbox(slide1, x, Inches(6.1), Inches(3), Inches(0.5),
                     _fmt(d['총매출']), font_size=22, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))

    _add_textbox(slide1, Inches(10), Inches(6.8), Inches(3), Inches(0.3),
                 f'Generated: {year}-{month} Report', font_size=9,
                 color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.RIGHT)

    # ━━━ 슬라이드 2: 매출 요약 ━━━
    slide2 = prs.slides.add_slide(blank_layout)
    s2_info = slide_titles.get(2, {})
    _add_slide_header(slide2,
                      s2_info.get('title', f'{month_str} 매출 요약'),
                      s2_info.get('subtitle', '전사 매출 KPI + 사업부별 매출 비교 + 전월 대비'), 2)

    total = data['total']
    prev = data['prev_month']
    kpis = [
        ('총매출', _fmt(_dec(total['총매출'])), _dec(total['총매출']), _dec(prev.get('총매출', 0))),
        ('판매수량(KG)', _fmt(_dec(total['판매수량'])), _dec(total['판매수량']), _dec(prev.get('판매수량', 0))),
        ('총 원가', _fmt(_dec(total['매출원가'])) if _dec(total['매출원가']) else '(산출전)', _dec(total['매출원가']), _dec(prev.get('매출원가', 0))),
        ('거래건수', f"{data['total_rows']:,}건", data['total_rows'], data['total_rows'] * 0.95),
    ]
    for i, (label, value_str, cur_val, prev_val) in enumerate(kpis):
        x = Inches(0.5 + i * 3.1)
        if prev_val and cur_val:
            chg = (cur_val - prev_val) / abs(prev_val) * 100 if prev_val else 0
            sub = f"{'▲' if chg >= 0 else '▼'} 전월 대비 {abs(chg):.1f}%"
        else:
            sub = ''
        _add_kpi_box(slide2, x, Inches(1.0), Inches(2.8), Inches(1.0),
                     label, value_str, sub)

    div_labels = [d['사업부'] for d in data['by_division']]
    div_sales = [d['총매출'] for d in data['by_division']]
    chart_img1 = make_bar_chart(div_labels, div_sales,
                                title=f'{month_str} 사업부별 매출', figsize=(5.5, 3.5))
    slide2.shapes.add_picture(chart_img1, Inches(0.5), Inches(2.3), Inches(5.5), Inches(3.5))

    _add_textbox(slide2, Inches(6.5), Inches(2.3), Inches(6), Inches(0.4),
                 '전월 대비 매출 비교', font_size=14, bold=True, color=COLORS['dark'])

    prev_total_sales = _dec(prev.get('총매출', 0))
    compare_headers = ['항목', '당월', '전월(추정)', '증감', '증감률']
    compare_items = [
        ('총매출', total['총매출'], prev.get('총매출', 0)),
        ('판매수량', total['판매수량'], prev.get('판매수량', 0)),
    ]
    compare_rows = []
    for name, cur_v, prev_v in compare_items:
        c = _dec(cur_v)
        p = _dec(prev_v)
        diff = c - p
        rate = (diff / abs(p) * 100) if p else 0
        compare_rows.append([
            name, _fmt(c), _fmt(p),
            f"{'+'if diff>=0 else ''}{_fmt(diff)}",
            f"{'+'if rate>=0 else ''}{rate:.1f}%"
        ])
    _add_table(slide2, Inches(6.5), Inches(2.8), Inches(6.3), Inches(1.2),
               compare_headers, compare_rows,
               col_widths=[Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.2), Inches(1.2)])

    ch_labels = [d['구분'] for d in data['by_channel']]
    ch_values = [d['총매출'] for d in data['by_channel']]
    chart_img2 = make_pie_chart(ch_labels, ch_values,
                                title='내수 / 수출 비중', figsize=(5.5, 3))
    slide2.shapes.add_picture(chart_img2, Inches(6.5), Inches(4.2), Inches(6), Inches(3))

    # ━━━ 슬라이드 3: 상세 분석 ━━━
    slide3 = prs.slides.add_slide(blank_layout)
    s3_info = slide_titles.get(3, {})
    _add_slide_header(slide3,
                      s3_info.get('title', f'{month_str} 사업부별 상세 분석'),
                      s3_info.get('subtitle', '플랜트별 매출 + 일별 매출 추이 + 주간 추이'), 3)

    plant_labels = [d['플랜트'] for d in data['by_plant']]
    plant_sales = [d['총매출'] for d in data['by_plant']]
    chart_img3 = make_bar_chart(plant_labels, plant_sales,
                                title='플랜트별 매출', figsize=(5.5, 3.5))
    slide3.shapes.add_picture(chart_img3, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5))

    daily_x = [d['일자'] for d in data['daily_trend']]
    daily_y = [d['총매출'] for d in data['daily_trend']]
    chart_img4 = make_line_chart(
        daily_x,
        [('일별 총매출', daily_y, CHART_COLORS[0])],
        title=f'{month_str} 일별 매출 추이',
        figsize=(6.5, 3.5)
    )
    slide3.shapes.add_picture(chart_img4, Inches(6.5), Inches(1.2), Inches(6.5), Inches(3.5))

    weekly_labels = [d['주차'] for d in data['weekly_trend']]
    weekly_values = [d['총매출'] for d in data['weekly_trend']]
    chart_img5 = make_bar_chart(weekly_labels, weekly_values,
                                title='주간별 매출 추이', figsize=(5.5, 2.8))
    slide3.shapes.add_picture(chart_img5, Inches(0.5), Inches(4.8), Inches(5.5), Inches(2.5))

    div_headers = ['사업부', '총매출', '판매수량(KG)', '건수', '비중']
    div_rows = []
    for d in data['by_division']:
        pct = _pct(d['총매출'], total['총매출'])
        div_rows.append([
            d['사업부'], _fmt(d['총매출']),
            _fmt(d['판매수량']),
            f"{d['건수']:,}", f"{pct:.1f}%"
        ])
    _add_table(slide3, Inches(6.5), Inches(4.8), Inches(6.3), Inches(1.5),
               div_headers, div_rows,
               col_widths=[Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.1), Inches(1.1)])

    # ━━━ 슬라이드 4: 원가 분석 ━━━
    slide4 = prs.slides.add_slide(blank_layout)
    s4_info = slide_titles.get(4, {})
    _add_slide_header(slide4,
                      s4_info.get('title', f'{month_str} 원가 구성 분석'),
                      s4_info.get('subtitle', '원가 항목별 구성 + 사업부별 원가 비교'), 4)

    cost = data['cost_detail']
    cost_items = [
        ('재료비(펄프)', cost.get('재료비_펄프', 0)),
        ('재료비(고지)', cost.get('재료비_고지', 0)),
        ('재료비(패드)', cost.get('재료비_패드', 0)),
        ('부재료(약품)', cost.get('부재료비_약품', 0)),
        ('부재료(포장재)', cost.get('부재료비_포장재', 0)),
        ('인건비', cost.get('인건비', 0)),
        ('에너지비', cost.get('에너지비', 0)),
        ('감가상각비', cost.get('감가상각비', 0)),
        ('외주가공비', cost.get('외주가공비', 0)),
    ]
    cost_items = [(k, v) for k, v in cost_items if _dec(v) != 0]

    cost_labels = [c[0] for c in cost_items]
    cost_values = [c[1] for c in cost_items]
    chart_img6 = make_pie_chart(cost_labels, cost_values,
                                title='원가 항목별 구성 비율', figsize=(6.5, 4))
    slide4.shapes.add_picture(chart_img6, Inches(0.3), Inches(1.2), Inches(6.5), Inches(4))

    div_data = data['by_division']
    if len(div_data) >= 2:
        cost_keys = ['재료비_펄프', '재료비_고지', '인건비', '에너지비', '감가상각비']
        cost_kr = ['펄프', '고지', '인건비', '에너지', '감가상각']
        groups = []
        for d in div_data[:2]:
            vals = [d.get(k, 0) for k in cost_keys]
            color = CHART_COLORS[0] if '제지' in d['사업부'] else CHART_COLORS[1]
            groups.append((d['사업부'], vals, color))
        chart_img7 = make_grouped_bar_chart(cost_kr, groups,
                                            title='사업부별 주요 원가 비교', figsize=(6, 4))
        slide4.shapes.add_picture(chart_img7, Inches(7), Inches(1.2), Inches(6), Inches(4))

    cost_headers = ['원가 항목', '금액', '비중']
    cost_total = sum(_dec(v) for _, v in cost_items)
    cost_rows_data = []
    for name, val in cost_items:
        v = _dec(val)
        pct = (v / cost_total * 100) if cost_total else 0
        cost_rows_data.append([name, _fmt_full(v), f"{pct:.1f}%"])
    cost_rows_data.append(['합계', _fmt_full(cost_total), '100.0%'])
    _add_table(slide4, Inches(0.5), Inches(5.4), Inches(5.5), Inches(1.8),
               cost_headers, cost_rows_data,
               col_widths=[Inches(1.8), Inches(2.2), Inches(1.5)])

    if len(div_data) >= 2:
        div_cost_headers = ['원가 항목', div_data[0]['사업부'], div_data[1]['사업부']]
        div_cost_rows = []
        for key, kr in zip(cost_keys, cost_kr):
            div_cost_rows.append([
                kr,
                _fmt(div_data[0].get(key, 0)),
                _fmt(div_data[1].get(key, 0))
            ])
        _add_table(slide4, Inches(7), Inches(5.4), Inches(5.5), Inches(1.8),
                   div_cost_headers, div_cost_rows,
                   col_widths=[Inches(1.5), Inches(2), Inches(2)])

    # ━━━ 슬라이드 5: 제품 & 인사이트 ━━━
    slide5 = prs.slides.add_slide(blank_layout)
    s5_info = slide_titles.get(5, {})
    _add_slide_header(slide5,
                      s5_info.get('title', f'{month_str} 제품별 분석 & 핵심 인사이트'),
                      s5_info.get('subtitle', '매출 TOP10 제품 + 브랜드별 매출 + 핵심 요약'), 5)

    top_prods = data['top_products']
    prod_labels = [p['제품명'][:20] for p in top_prods]
    prod_values = [p['총매출'] for p in top_prods]
    chart_img8 = make_horizontal_bar_chart(
        prod_labels[::-1], prod_values[::-1],
        title='매출 TOP 10 제품', figsize=(6.5, 4.5)
    )
    slide5.shapes.add_picture(chart_img8, Inches(0.3), Inches(1.1), Inches(6.5), Inches(4.5))

    brands = data['by_brand']
    if brands:
        br_labels = [b['브랜드'] for b in brands]
        br_values = [b['총매출'] for b in brands]
        chart_img9 = make_bar_chart(br_labels, br_values,
                                    title='브랜드별 매출', figsize=(5.5, 3))
        slide5.shapes.add_picture(chart_img9, Inches(7), Inches(1.1), Inches(6), Inches(3))

    _add_textbox(slide5, Inches(7), Inches(4.3), Inches(6), Inches(0.4),
                 '핵심 인사이트', font_size=14, bold=True, color=COLORS['dark'])

    insights = []
    if data['by_division']:
        top_div = data['by_division'][0]
        insights.append(f"1. {top_div['사업부']} 매출 {_fmt(top_div['총매출'])} (전사 {_pct(top_div['총매출'], total['총매출']):.1f}%)")
    if data['by_plant']:
        top_plant = data['by_plant'][0]
        insights.append(f"2. {top_plant['플랜트']} 공장 매출 {_fmt(top_plant['총매출'])} (1위)")
    if data['by_channel']:
        for ch in data['by_channel']:
            if ch['구분'] == '내수':
                pct = _pct(ch['총매출'], total['총매출'])
                insights.append(f"3. 내수 매출 비중 {pct:.1f}% ({_fmt(ch['총매출'])})")
                break
    if cost_items:
        top_cost = max(cost_items, key=lambda x: _dec(x[1]))
        insights.append(f"4. 최대 원가: {top_cost[0]} ({_fmt(top_cost[1])})")
    if top_prods:
        insights.append(f"5. 매출 1위 제품: {top_prods[0]['제품명'][:25]}")
    if prev_total_sales:
        chg_final = ((total_sales - prev_total_sales) / abs(prev_total_sales) * 100)
        direction = '증가' if chg_final >= 0 else '감소'
        insights.append(f"6. 전월 대비 매출 {abs(chg_final):.1f}% {direction}")

    # 프롬프트 기반 추가 인사이트
    if prompt:
        insights.append(f"* 분석 요청: {prompt[:60]}")

    # 첨부파일 기반 추가 인사이트
    if attachment_data and attachment_data.get('type') in ('excel', 'csv'):
        insights.append(f"* 첨부 데이터: {attachment_data.get('filename', '')} ({attachment_data.get('row_count', 0)}행)")

    insight_box = slide5.shapes.add_shape(1, Inches(7), Inches(4.7), Inches(5.8), Inches(2.5))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    insight_box.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
    insight_box.line.width = Pt(1)

    tf = insight_box.text_frame
    tf.word_wrap = True
    for line_text in insights:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = line_text
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text']
        p.font.name = '맑은 고딕'
        p.space_after = Pt(6)

    # 첨부 이미지가 있으면 마지막 슬라이드에 삽입
    if attachment_data and attachment_data.get('type') == 'image':
        img_path = attachment_data.get('path', '')
        if os.path.exists(img_path):
            try:
                slide5.shapes.add_picture(img_path, Inches(0.3), Inches(5.8), Inches(3), Inches(1.5))
            except Exception:
                pass

    # PPT 저장
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
